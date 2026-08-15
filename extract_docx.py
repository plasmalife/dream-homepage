from pptx import Presentation
import io

prs = Presentation(r'e:\ANTI_gravity\dream homepage\쿠팡상품\강타자\강타자_상품_상세설명서.pptx')
output = io.StringIO()
for i, slide in enumerate(prs.slides):
    output.write(f'=== Slide {i+1} ===\n')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    output.write(text + '\n')
    output.write('\n')

with open(r'e:\ANTI_gravity\dream homepage\수정사항\gangtaja_pptx.txt', 'w', encoding='utf-8') as f:
    f.write(output.getvalue())
print('Done')
